#!/usr/bin/env python3
"""Build the JDBC live-session deck (import into Google Slides)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "PayNest-JDBC-H2.pptx"

NAVY = RGBColor(0x0F, 0x2A, 0x44)
TEAL = RGBColor(0x1F, 0x7A, 0x6B)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x64, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF4, 0xF7, 0xF8)
CARD = RGBColor(0xE8, 0xF1, 0xEF)
WARM = RGBColor(0xFF, 0xF3, 0xE0)


def set_run(run, text, *, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, lines, *, size=18, bold=False, color=INK, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        set_run(r, line, size=size, bold=bold, color=color, font=font)
    return box


def add_bg(slide, color=CREAM):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def title_block(slide, kicker, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.35), [kicker.upper()], size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(11.8), Inches(0.9), [title], size=32, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.75), Inches(11.8), Inches(0.55), [subtitle], size=16, color=MUTED)


def card(slide, left, top, width, height, lines, *, fill=CARD, size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        set_run(r, line, size=size, color=INK)
        p.space_after = Pt(6)


def code_card(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        set_run(r, line, size=14, color=WHITE, font="Consolas")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    footer = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = TEAL
    footer.line.fill.background()
    add_textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.4), ["PayNest live session"], size=16, bold=True, color=TEAL)
    add_textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.9), ["JDBC with H2"], size=44, bold=True, color=WHITE)
    add_textbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.5), ["Connect · Create a table · Insert · Query"], size=20, color=WHITE)
    add_textbox(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.4), ["23 July 2026"], size=14, color=RGBColor(0xB8, 0xC7, 0xD1))

    # 2 Goals
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Today", "What you will learn", "Four practical database steps in Java.")
    card(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(12),
        Inches(3.8),
        [
            "1. Open a JDBC connection to H2",
            "2. Create a table with CREATE TABLE",
            "3. Insert rows with PreparedStatement",
            "4. Query rows with SELECT + ResultSet",
            "5. Inspect the same data in DBeaver",
        ],
    )

    # 3 H2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Database", "What is H2?", "A small relational database that runs with your Java app.")
    card(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(5.8),
        Inches(3.8),
        [
            "Stores tables and rows like other SQL databases",
            "Can keep data in a local file",
            "Already on the PayNest Maven classpath",
            "Ideal for learning JDBC quickly",
        ],
    )
    card(
        s,
        Inches(6.8),
        Inches(2.6),
        Inches(5.8),
        Inches(3.8),
        [
            "Tools for this session",
            "• DBeaver to browse tables",
            "• Java DriverManager to connect",
            "• Not pgAdmin (Postgres only)",
        ],
        fill=WARM,
    )

    # 4 Files
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "On disk", "Two files under data/", "Only one of them is the real database.")
    card(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(5.8),
        Inches(3.8),
        [
            "paynest.mv.db",
            "• The real database",
            "• Tables and rows live here",
            "• Connect DBeaver to this DB",
        ],
    )
    card(
        s,
        Inches(6.8),
        Inches(2.6),
        Inches(5.8),
        Inches(3.8),
        [
            "paynest.trace.db",
            "• SQL debug / trace log",
            "• Created with TRACE_LEVEL_FILE",
            "• Do not open this as the database",
        ],
        fill=WARM,
    )

    # 5 Connect
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Step 1", "Create a database connection")
    code_card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(12),
        Inches(2.2),
        'String url = "jdbc:h2:file:./data/paynest;AUTO_SERVER=TRUE;TRACE_LEVEL_FILE=3";\n'
        "Connection connection = DriverManager.getConnection(url);",
    )
    card(
        s,
        Inches(0.7),
        Inches(5.0),
        Inches(12),
        Inches(1.8),
        [
            "file: = durable on disk",
            "AUTO_SERVER=TRUE = share the file with DBeaver",
            "TRACE_LEVEL_FILE=3 = write data/paynest.trace.db",
        ],
    )

    # 6 Create
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Step 2", "Create a table")
    code_card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(12),
        Inches(3.8),
        "CREATE TABLE IF NOT EXISTS products (\n"
        "  id    INT PRIMARY KEY,\n"
        "  name  VARCHAR(100) NOT NULL,\n"
        "  price DOUBLE NOT NULL\n"
        ");",
    )

    # 7 Insert
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Step 3", "Insert data safely")
    code_card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(12),
        Inches(3.2),
        "PreparedStatement insert = connection.prepareStatement(\n"
        '  "INSERT INTO products (id, name, price) VALUES (?, ?, ?)");\n'
        "insert.setInt(1, 1);\n"
        'insert.setString(2, "Laptop");\n'
        "insert.setDouble(3, 12000);\n"
        "insert.executeUpdate();",
    )
    card(
        s,
        Inches(0.7),
        Inches(5.9),
        Inches(12),
        Inches(1.0),
        ["Use ? placeholders — never paste raw values into SQL strings."],
        size=15,
    )

    # 8 Query
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Step 4", "Query the data")
    code_card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(12),
        Inches(3.8),
        "ResultSet rows = statement.executeQuery(\n"
        '  "SELECT id, name, price FROM products ORDER BY id");\n'
        "while (rows.next()) {\n"
        '  int id = rows.getInt("id");\n'
        '  String name = rows.getString("name");\n'
        '  double price = rows.getDouble("price");\n'
        "}",
    )

    # 9 Code location
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Code", "One place to look", "All four steps are in PayNestApplication.")
    card(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(12),
        Inches(3.8),
        [
            "File: src/main/java/com/paynestsystem/app/PayNestApplication.java",
            "Method: runJdbcDemo(...)",
            "Run: mvn -q compile exec:java",
            "Then in DBeaver: SELECT * FROM products;",
        ],
    )

    # 10 DBeaver
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "GUI", "Inspect with DBeaver")
    card(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(12),
        Inches(3.8),
        [
            "1. New connection → H2",
            "2. URL: jdbc:h2:file:/ABS/PATH/PayNest/data/paynest;AUTO_SERVER=TRUE",
            "3. Test connection → Finish",
            "4. Run: SELECT * FROM products ORDER BY id;",
            "Tip: run the Java app once first so data/ and the table exist",
        ],
    )

    # 11 Recap
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_block(s, "Recap", "Remember the loop")
    steps = [
        ("1", "Connect", "DriverManager"),
        ("2", "Create", "CREATE TABLE"),
        ("3", "Insert", "PreparedStatement"),
        ("4", "Query", "ResultSet"),
    ]
    left = 0.7
    for num, label, api in steps:
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.7), Inches(2.9), Inches(3.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, num, size=36, bold=True, color=TEAL)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, label, size=20, bold=True, color=NAVY)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, api, size=13, color=MUTED)
        left += 3.1

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
