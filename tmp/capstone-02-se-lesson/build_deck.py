#!/usr/bin/env python3
"""Build Capstone 2 software-engineering live lesson deck (Google Slides–compatible PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
DOCS_OUT = (
    Path(__file__).resolve().parents[2]
    / "docs"
    / "live-sessions"
    / "capstone-02-design-to-interfaces"
    / "capstone-02-se-lesson.pptx"
)
OUT_FILE = OUT_DIR / "capstone-02-se-lesson.pptx"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED_SOFT = RGBColor(0xB9, 0x1C, 0x1C)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
ROW_ALT = RGBColor(0xF1, 0xF5, 0xF9)

FONT_TITLE = "Poppins"
FONT_BODY = "Lato"
FONT_CODE = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box, tf


def add_multiline(tf, lines, size=14, color=OFF_WHITE, font=FONT_CODE, spacing=2):
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(spacing)


def add_bullets(tf, items, size=16, color=SLATE, spacing=6):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)


def add_header_bar(slide, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL)
    if kicker:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(0.32),
            Inches(12),
            Inches(0.35),
            kicker.upper(),
            size=11,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
        )


def add_footer(slide, text="PayNest Capstone Programme · Capstone 2 SE Lesson"):
    add_textbox(
        slide, Inches(0.7), Inches(7.05), Inches(11.5), Inches(0.3), text, size=9, color=SLATE
    )


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_title(slide, title, subtitle=None):
    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.75),
        Inches(12),
        Inches(0.7),
        title,
        size=30,
        bold=True,
        color=NAVY,
        font=FONT_TITLE,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(1.4),
            Inches(12),
            Inches(0.45),
            subtitle,
            size=16,
            color=SLATE,
        )


def section_divider(prs, block, title, time_label, objective):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.8),
        Inches(11),
        Inches(0.4),
        f"BLOCK {block}  ·  {time_label}",
        size=14,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(2.4),
        Inches(11.5),
        Inches(1.4),
        title,
        size=36,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(4.2),
        Inches(11),
        Inches(1.2),
        f"Learning objective: {objective}",
        size=18,
        color=OFF_WHITE,
    )
    notes(
        slide,
        f"SECTION OPENER — Block {block}\n"
        f"Objective: {objective}\n"
        "Pause. Tell students what they will be able to explain after this block.\n"
        "Do not dive into syntax yet.",
    )
    return slide


def exercise_slide(prs, title, prompt, minutes="5 min"):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Student exercise")
    content_title(slide, title)
    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.4), WHITE)
    add_rect(slide, Inches(0.7), Inches(1.95), Inches(0.12), Inches(4.4), AMBER)
    add_textbox(
        slide,
        Inches(1.1),
        Inches(2.2),
        Inches(3),
        Inches(0.4),
        minutes.upper(),
        size=13,
        bold=True,
        color=AMBER,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(1.1),
        Inches(2.75),
        Inches(11),
        Inches(3.2),
        prompt,
        size=20,
        color=NAVY,
    )
    add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Block 0 — Opening + Capstone 1 recap
# ---------------------------------------------------------------------------


def slide_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    add_round_rect(slide, Inches(9.2), Inches(1.2), Inches(3.6), Inches(5.1), TEAL_LIGHT)
    add_round_rect(slide, Inches(9.6), Inches(1.6), Inches(2.8), Inches(4.3), TEAL)

    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(7.8),
        Inches(0.4),
        "CAPSTONE 2  ·  LIVE LESSON",
        size=14,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(2.1),
        Inches(7.8),
        Inches(1.8),
        "From Design Pain\nto Interfaces",
        size=40,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(4.2),
        Inches(7.5),
        Inches(1.0),
        "A software engineering class — not a Java syntax class.\n2–3 hours · live walkthrough of the PayNest solution.",
        size=18,
        color=OFF_WHITE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(6.2),
        Inches(7),
        Inches(0.4),
        "PayNest · Second-year Computer Science",
        size=13,
        color=SLATE,
    )
    notes(
        slide,
        "OPENING\n"
        "Objective: Set the frame — today we learn why Capstone 2's design exists.\n"
        "Why this exists: Students who jump to 'implements PaymentMethod' miss the engineering reason.\n"
        "Analogy: Learning traffic lights by memorising colours vs understanding why intersections need them.\n"
        "Live coding: None yet. Open PayNestApplication and leave Capstone 1 summary on screen.\n"
        "Exercise: None.\n"
        "Common mistakes: Starting with 'today we learn interfaces.'\n"
        "Ask: What do you already trust about Capstone 1 orders?\n"
        "Transition: We start from what you already built.",
    )


def slide_agenda(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Agenda")
    content_title(slide, "Where we are going", "Design problem first. Syntax last.")

    rows = [
        ("0", "15m", "Capstone 1 recap — objects you already own"),
        ("1", "25m", "Why software changes — the if/else trap"),
        ("2", "20m", "Responsibility, cohesion, coupling, composition"),
        ("3", "25m", "Abstractions & the PaymentMethod contract"),
        ("4", "30m", "Polymorphism, OCP, DIP"),
        ("5", "20m", "Processor, checkout, separation of concerns"),
        ("6", "20m", "Business rules, DRY, testing"),
        ("7", "15m", "Map everything to the Capstone 2 rubric"),
    ]
    for i, (num, time, label) in enumerate(rows):
        top = Inches(1.95 + i * 0.58)
        add_round_rect(slide, Inches(0.7), top, Inches(0.55), Inches(0.48), TEAL)
        add_textbox(
            slide,
            Inches(0.7),
            top + Inches(0.05),
            Inches(0.55),
            Inches(0.4),
            num,
            size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            Inches(1.4),
            top + Inches(0.08),
            Inches(1.0),
            Inches(0.35),
            time,
            size=13,
            bold=True,
            color=TEAL,
        )
        add_textbox(
            slide, Inches(2.5), top + Inches(0.08), Inches(10), Inches(0.4), label, size=16, color=NAVY
        )
    add_footer(slide)
    notes(
        slide,
        "AGENDA\n"
        "Tell students: interfaces appear only after we feel the pain of change.\n"
        "Breaks are embedded; total ~150–180 minutes.",
    )


def slide_not_syntax(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Framing")
    content_title(slide, "This is not 'today we learn interfaces'")

    left = [
        ("We will NOT", "Start with interface keyword quizzes\nMemorise SOLID acronyms first\nTreat Capstone 2 as busywork"),
        ("We WILL", "Start from a business change request\nFeel why tightly coupled code hurts\nDiscover abstractions as a design response"),
    ]
    for i, (head, body) in enumerate(left):
        x = Inches(0.7 + i * 6.2)
        bg = NAVY if i == 0 else TEAL
        add_round_rect(slide, x, Inches(2.0), Inches(5.9), Inches(4.2), bg)
        add_textbox(
            slide,
            x + Inches(0.4),
            Inches(2.3),
            Inches(5.1),
            Inches(0.5),
            head,
            size=20,
            bold=True,
            color=TEAL_LIGHT if i == 0 else WHITE,
            font=FONT_TITLE,
        )
        add_textbox(
            slide,
            x + Inches(0.4),
            Inches(3.1),
            Inches(5.1),
            Inches(2.8),
            body,
            size=18,
            color=OFF_WHITE,
        )
    add_footer(slide)
    notes(
        slide,
        "FRAMING\n"
        "Objective: Reset student expectations away from syntax drills.\n"
        "Why: Capstone 2 exists because Ops asked for extensible checkout — not because 'interfaces are in the curriculum.'\n"
        "Analogy: Doctors don't start with 'today we learn scalpels' — they start with the patient problem.\n"
        "Ask: When has a feature request forced you to rewrite working code?\n"
        "Transition: Let's recall Capstone 1 — the code that must stay trustworthy.",
    )


def slide_you_already_know(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 0 · Recap")
    content_title(slide, "You already know these tools", "Capstone 2 reuses them under new pressure.")

    cards = [
        ("Classes & objects", "Blueprints and instances"),
        ("Constructors", "Valid starting state"),
        ("Encapsulation", "Private fields, controlled access"),
        ("Getters", "Safe read of state"),
        ("Lists", "Order lines collection"),
        ("Totals", "OrderItem + Order math"),
        ("Maven", "Build & run"),
        ("JUnit", "Prove behaviour"),
    ]
    for i, (title, sub) in enumerate(cards):
        col, row = i % 4, i // 4
        left = Inches(0.7 + col * 3.1)
        top = Inches(2.1 + row * 2.2)
        add_round_rect(slide, left, top, Inches(2.95), Inches(1.9), WHITE)
        add_rect(slide, left, top, Inches(2.95), Inches(0.08), TEAL)
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.35),
            Inches(2.5),
            Inches(0.7),
            title,
            size=15,
            bold=True,
            color=NAVY,
            font=FONT_TITLE,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.1),
            Inches(2.5),
            Inches(0.5),
            sub,
            size=13,
            color=SLATE,
        )
    add_footer(slide)
    notes(
        slide,
        "RECAP TOOLS\n"
        "Objective: Activate prior knowledge so Capstone 2 feels continuous.\n"
        "Why: New abstractions land on familiar OOP skills.\n"
        "Live coding: Flash Order.java fields and calculateTotal — no edits.\n"
        "Ask: Which of these protect the R12,400 demo total?\n"
        "Transition: Open the Capstone 1 path in PayNestApplication.",
    )


def slide_c1_walk(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 0 · Live walk")
    content_title(slide, "Capstone 1 in one breath", "Open PayNestApplication — lines building the order desk.")

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(7.4), Inches(4.5), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.25), Inches(6.9), Inches(4.0), "")
    add_multiline(
        tf,
        [
            "Product laptop = new Product(1, \"Laptop\", 12000);",
            "Product mouse  = new Product(2, \"Mouse\", 200);",
            "Customer c = new Customer(1, \"John Smith\", ...);",
            "Order order = orderService.createOrder(1, c);",
            "orderService.addProductsToOrder(order, laptop, 1);",
            "orderService.addProductsToOrder(order, mouse, 2);",
            "order.printSummary();  // Total: R12400",
        ],
        size=13,
    )

    add_round_rect(slide, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.5), WHITE)
    add_textbox(
        slide,
        Inches(8.7),
        Inches(2.25),
        Inches(3.8),
        Inches(0.4),
        "Trust checklist",
        size=15,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(8.7), Inches(2.85), Inches(3.8), Inches(3.4), "")
    add_bullets(
        tf2,
        [
            "Objects hold real data",
            "Encapsulation hides lists",
            "One path for grand total",
            "Receipt matches math",
        ],
        size=14,
    )
    add_footer(slide)
    notes(
        slide,
        "C1 WALK\n"
        "Objective: Re-establish the commerce kernel as trustworthy ground truth.\n"
        "Why: Checkout must charge that same total — no new arithmetic story.\n"
        "Analogy: The till receipt must match the shelf prices.\n"
        "Live coding: Run mvn exec:java (or highlight printSummary output). Stop before payment if showing C1-only narrative first.\n"
        "Ask: If Ops asks for three payment rails, what must NOT change?\n"
        "Transition: Hook — 'Ops wants card, EFT, and wallet. What breaks if we hard-code that?'",
    )


def slide_hook(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(2.4),
        Inches(11.5),
        Inches(1.5),
        "Ops wants card, EFT, and wallet.\nWhat breaks if we hard-code that?",
        size=32,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1.5),
        Inches(4.5),
        Inches(10.3),
        Inches(0.8),
        "Hold that question. We will design the wrong way first — on purpose.",
        size=18,
        color=TEAL_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    notes(
        slide,
        "HOOK\n"
        "Pause for 20 seconds of silence after the question.\n"
        "Collect 2–3 student guesses on the whiteboard.\n"
        "Transition into Block 1.",
    )


# ---------------------------------------------------------------------------
# Block 1 — Change & tight coupling
# ---------------------------------------------------------------------------


def slide_why_software_changes(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 1 · Change")
    content_title(
        slide,
        "Why software changes",
        "Working code is not finished code — the business keeps asking.",
    )

    cards = [
        ("New channels", "Card today. Wallet tomorrow. BNPL next sprint."),
        ("Ops pressure", "Three copy-pasted checkout paths = three bug farms."),
        ("Same money", "Product pitch: same order total; only the rail changes."),
    ]
    for i, (h, b) in enumerate(cards):
        left = Inches(0.7 + i * 4.1)
        add_round_rect(slide, left, Inches(2.2), Inches(3.9), Inches(3.5), WHITE)
        add_rect(slide, left, Inches(2.2), Inches(3.9), Inches(0.08), TEAL)
        add_textbox(
            slide,
            left + Inches(0.3),
            Inches(2.6),
            Inches(3.3),
            Inches(0.8),
            h,
            size=18,
            bold=True,
            color=NAVY,
            font=FONT_TITLE,
        )
        add_textbox(
            slide, left + Inches(0.3), Inches(3.6), Inches(3.3), Inches(1.6), b, size=15, color=SLATE
        )
    add_footer(slide)
    notes(
        slide,
        "WHY SOFTWARE CHANGES\n"
        "Objective: Students can explain change pressure as a design force.\n"
        "Why this concept exists: SE is managing inevitable change without destroying working systems.\n"
        "Analogy: A restaurant menu grows — kitchen layout must absorb new dishes without rebuilding the building.\n"
        "Live coding: Read Capstone 2 brief scenario aloud (merchants loved C1; now want checkout).\n"
        "Exercise: None yet.\n"
        "Common mistakes: Treating requirements as frozen forever.\n"
        "Ask: What changed between Capstone 1 and Capstone 2 from Ops' point of view?\n"
        "Transition: Let's write the 'obvious' checkout — and feel why it fails.",
    )


def slide_bad_checkout(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 1 · Anti-pattern")
    content_title(slide, "The tempting design", "One method. Three rails. Plenty of if/else.")

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(8.0), Inches(4.6), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.2), Inches(7.5), Inches(4.2), "")
    add_multiline(
        tf,
        [
            "// DO NOT ship this — teaching anti-pattern",
            "void checkout(String type) {",
            "  double total = calculateTotal();",
            "  if (type.equals(\"CARD\")) {",
            "    System.out.println(\"Charging card...\");",
            "    // card-specific logic...",
            "  } else if (type.equals(\"EFT\")) {",
            "    // eft-specific logic...",
            "  } else if (type.equals(\"WALLET\")) {",
            "    // wallet-specific logic...",
            "  }",
            "}",
        ],
        size=13,
    )

    add_round_rect(slide, Inches(9.0), Inches(1.95), Inches(3.7), Inches(4.6), WHITE)
    add_textbox(
        slide,
        Inches(9.25),
        Inches(2.2),
        Inches(3.3),
        Inches(0.4),
        "Smell test",
        size=15,
        bold=True,
        color=RED_SOFT,
        font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(9.25), Inches(2.8), Inches(3.3), Inches(3.4), "")
    add_bullets(
        tf2,
        [
            "Order knows every rail",
            "New rail = edit Order",
            "Strings are fragile",
            "Messaging duplicated",
        ],
        size=14,
        color=SLATE,
    )
    add_footer(slide)
    notes(
        slide,
        "ANTI-PATTERN CHECKOUT\n"
        "Objective: Recognise tight coupling between Order and payment types.\n"
        "Why: Illustrating failure modes teaches better than praising patterns.\n"
        "Analogy: One cashier who must know every bank's phone number by heart.\n"
        "Live coding: Type this anti-pattern on a scratch slide or scratch class — do NOT put it in Order.java.\n"
        "Common mistakes: Thinking 'just one more else if' is fine forever.\n"
        "Ask: How many places change when Ops adds Buy Now Pay Later?\n"
        "Transition: Let's do the BNPL exercise.",
    )


def slide_exercise_bnpl_pain(prs):
    slide = exercise_slide(
        prs,
        "Exercise: add Buy Now Pay Later",
        "Using the if/else checkout, list every line or place you must touch to add "
        "BuyNowPayLater.\n\n"
        "Include: Order method, messages, possible typos in string labels, tests, demo wiring.\n\n"
        "Be specific. Write the list.",
        "4–5 min",
    )
    notes(
        slide,
        "EXERCISE — BNPL PAIN\n"
        "Objective: Make change cost visceral.\n"
        "Collect lists on the board. Expect: else-if branch, string constant, print messages, "
        "maybe a switch elsewhere, demo main, tests.\n"
        "Debrief: This list IS the maintainability problem.\n"
        "Transition: Name the SE concepts for what we just felt.",
    )


def slide_tight_coupling(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 1 · Coupling")
    content_title(
        slide,
        "Tight coupling hurts maintainability",
        "When A cannot change without rewriting B, they are glued together.",
    )

    pairs = [
        ("Tight", "Order hard-codes CARD / EFT / WALLET branches"),
        ("Loose (goal)", "Order asks any payment rail to process(amount)"),
        ("Symptom", "Every new rail reopens a 'finished' class"),
        ("Risk", "Regression in totals while editing payment noise"),
    ]
    for i, (h, b) in enumerate(pairs):
        top = Inches(2.05 + i * 1.1)
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(0.95), WHITE)
        add_textbox(
            slide,
            Inches(1.0),
            top + Inches(0.25),
            Inches(2.2),
            Inches(0.45),
            h,
            size=16,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
        )
        add_textbox(
            slide, Inches(3.4), top + Inches(0.25), Inches(8.8), Inches(0.5), b, size=16, color=NAVY
        )
    add_footer(slide)
    notes(
        slide,
        "TIGHT COUPLING\n"
        "Objective: Define coupling and why Capstone 2 attacks it.\n"
        "Why: Maintainability is a first-class SE goal for growing systems.\n"
        "Analogy: Headphones glued into a phone — upgrade either part and both break.\n"
        "Common mistakes: Confusing 'works in demo' with 'safe to extend'.\n"
        "Ask: Is Order's job 'own line items and totals' or 'know Visa vs EFT protocols'?\n"
        "Transition: If Order owns too much, we need vocabulary for good boundaries — Block 2.",
    )


# ---------------------------------------------------------------------------
# Block 2 — SE vocabulary from C1
# ---------------------------------------------------------------------------


def slide_responsibility(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 2 · Responsibility")
    content_title(
        slide,
        "Responsibility: what is Order for?",
        "Capstone 1 already answered — protect that answer.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.3), TEAL)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.3),
        Inches(5.2),
        Inches(0.5),
        "Order SHOULD",
        size=18,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(1.0), Inches(3.0), Inches(5.2), Inches(3.0), "")
    add_bullets(
        tf,
        [
            "Own its line items",
            "Validate adds via OrderItem",
            "Calculate the grand total",
            "Print a reconcilable summary",
        ],
        size=16,
        color=OFF_WHITE,
    )

    add_round_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.3), NAVY)
    add_textbox(
        slide,
        Inches(7.1),
        Inches(2.3),
        Inches(5.2),
        Inches(0.5),
        "Order should NOT",
        size=18,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(7.1), Inches(3.0), Inches(5.2), Inches(3.0), "")
    add_bullets(
        tf2,
        [
            "Know card network quirks",
            "Hard-code rail labels",
            "Duplicate total math",
            "Become a payment god class",
        ],
        size=16,
        color=OFF_WHITE,
    )
    add_footer(slide)
    notes(
        slide,
        "RESPONSIBILITY\n"
        "Objective: State Order's responsibility in one sentence.\n"
        "Why: Single Responsibility Principle emerges from clear ownership.\n"
        "Analogy: Kitchen — chef cooks; cashier takes payment. Mixing roles creates chaos at rush hour.\n"
        "Live coding: Open Order.java — point at addItem, calculateTotal, printSummary, getItems.\n"
        "Ask: Where does payment fit without polluting those duties?\n"
        "Transition: Cohesion and coupling name how well we grouped those duties.",
    )


def slide_cohesion_coupling(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 2 · Cohesion & coupling")
    content_title(slide, "Cohesion up. Coupling down.")

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(5.9), Inches(4.4), WHITE)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.2),
        Inches(5.3),
        Inches(0.5),
        "High cohesion",
        size=20,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.9),
        Inches(5.3),
        Inches(3.0),
        "Methods in a class work on the same job.\n\n"
        "OrderItem: quantity × product price.\n"
        "Order: collection of lines + grand total.\n\n"
        "Everything in the class pulls in the same direction.",
        size=16,
        color=NAVY,
    )

    add_round_rect(slide, Inches(6.9), Inches(1.95), Inches(5.9), Inches(4.4), WHITE)
    add_textbox(
        slide,
        Inches(7.2),
        Inches(2.2),
        Inches(5.3),
        Inches(0.5),
        "Low coupling",
        size=20,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(7.2),
        Inches(2.9),
        Inches(5.3),
        Inches(3.0),
        "Classes depend on as little as possible.\n\n"
        "Prefer: depend on a small contract.\n"
        "Avoid: depend on every concrete rail.\n\n"
        "Change stays local when coupling is low.",
        size=16,
        color=NAVY,
    )
    add_footer(slide)
    notes(
        slide,
        "COHESION & COUPLING\n"
        "Objective: Define both and spot them in Capstone 1 code.\n"
        "Why: These words let teams argue about design without personal taste.\n"
        "Analogy: A toolbox with only screwdrivers (cohesive) vs a junk drawer (low cohesion).\n"
        "Common mistakes: Giant Utils class — low cohesion, high coupling magnet.\n"
        "Ask: Is the if/else checkout high or low cohesion inside Order?\n"
        "Transition: Composition is how Order already builds complexity safely.",
    )


def slide_composition(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 2 · Composition")
    content_title(
        slide,
        "Composition: Order has OrderItems",
        "We build systems by combining objects — not by inheriting everything.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(7.5), Inches(4.3), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.3), Inches(7.0), Inches(3.8), "")
    add_multiline(
        tf,
        [
            "public class Order {",
            "  private final Customer customer;     // has-a",
            "  private final List<OrderItem> items; // has-many",
            "",
            "  public void addItem(Product p, int qty) {",
            "    items.add(new OrderItem(p, qty));",
            "  }",
            "",
            "  public List<OrderItem> getItems() {",
            "    return Collections.unmodifiableList(items);",
            "  }",
            "}",
        ],
        size=14,
    )

    add_round_rect(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.3), WHITE)
    _, tf2 = add_textbox(slide, Inches(8.8), Inches(2.3), Inches(3.7), Inches(3.8), "")
    add_bullets(
        tf2,
        [
            "Customer referenced, not subclassed",
            "Items owned privately",
            "Callers cannot clear() the list",
            "Checkout should compose a payment capability the same way",
        ],
        size=14,
    )
    add_footer(slide)
    notes(
        slide,
        "COMPOSITION\n"
        "Objective: Explain has-a vs inherits-a using Order.\n"
        "Why: Capstone 2 checkout should compose a PaymentMethod, not inherit CardPayment.\n"
        "Analogy: A playlist has songs; it is not a kind of song.\n"
        "Live coding: Highlight unmodifiable getItems — encapsulation protecting totals.\n"
        "Common mistakes: Public ArrayList field; recalculating totals in three places.\n"
        "Ask: Should Order extend CardPayment? Why does that sound wrong?\n"
        "Transition: We need a payment capability that Order can compose — that needs an abstraction.",
    )


def slide_exercise_vocab(prs):
    slide = exercise_slide(
        prs,
        "Exercise: name the smell",
        "In pairs, label the anti-pattern checkout with these words:\n"
        "responsibility · cohesion · coupling · composition\n\n"
        "One sentence each. Be ready to share.",
        "4 min",
    )
    notes(
        slide,
        "EXERCISE — VOCAB\n"
        "Sample answers: low cohesion (payment mixed with order math); high coupling to string rails; "
        "Order taking payment responsibility it shouldn't; composition missing (no payment object).\n"
        "Transition: Block 3 — invent the shared contract.",
    )


# ---------------------------------------------------------------------------
# Block 3 — Abstractions & interfaces
# ---------------------------------------------------------------------------


def slide_need_abstraction(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 3 · Abstraction")
    content_title(
        slide,
        "What do all payment rails share?",
        "Abstract thinking: ignore differences; keep the essential behaviour.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.1), Inches(11.9), Inches(2.0), WHITE)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.4),
        Inches(11.3),
        Inches(1.4),
        "Every rail must: (1) accept an amount in Rand and attempt to charge it,\n"
        "and (2) identify itself with a stable human-readable label for receipts and logs.",
        size=20,
        color=NAVY,
    )

    for i, label in enumerate(["Card", "EFT", "Wallet", "BNPL?"]):
        left = Inches(0.7 + i * 3.1)
        add_round_rect(slide, left, Inches(4.5), Inches(2.95), Inches(1.5), NAVY if i < 3 else TEAL)
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(4.9),
            Inches(2.5),
            Inches(0.7),
            label,
            size=20,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font=FONT_TITLE,
        )
    add_footer(slide)
    notes(
        slide,
        "NEED FOR ABSTRACTION\n"
        "Objective: Derive the shared behaviour before naming Java interfaces.\n"
        "Why abstractions exist: Call sites need a stable vocabulary while details vary.\n"
        "Analogy: 'Pay the bill' means the same whether you use cash, card, or transfer — "
        "the waiter doesn't run different restaurants for each.\n"
        "Live coding: Ask students to shout the two shared operations before opening PaymentMethod.java.\n"
        "Common mistakes: Jumping to 'interface' keyword before agreeing on behaviour.\n"
        "Ask: What must stay the same when the rail changes?\n"
        "Transition: That shared behaviour is a contract — in Java, an interface.",
    )


def slide_exercise_signatures(prs):
    slide = exercise_slide(
        prs,
        "Exercise: write the contract",
        "Before we open the file, write two method signatures that every payment rail must provide.\n\n"
        "Hint: charge an amount · describe yourself.\n\n"
        "Do not peek at PaymentMethod.java yet.",
        "3–4 min",
    )
    notes(
        slide,
        "EXERCISE — SIGNATURES\n"
        "Accept processPayment(double) / pay(double) / charge(double) and getType() / getName().\n"
        "Then reveal the real interface — celebrate near matches.\n"
        "Transition: Open PaymentMethod.java.",
    )


def slide_payment_method(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 3 · Interface")
    content_title(
        slide,
        "PaymentMethod — the contract",
        "An interface is a promise: any implementer can be used where this type is required.",
    )

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(7.6), Inches(4.5), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.25), Inches(7.1), Inches(4.0), "")
    add_multiline(
        tf,
        [
            "public interface PaymentMethod {",
            "  boolean processPayment(double amount);",
            "  String getPaymentType();",
            "}",
        ],
        size=18,
    )

    add_round_rect(slide, Inches(8.6), Inches(1.95), Inches(4.1), Inches(4.5), WHITE)
    _, tf2 = add_textbox(slide, Inches(8.9), Inches(2.3), Inches(3.6), Inches(3.9), "")
    add_bullets(
        tf2,
        [
            "No rail details here",
            "No if/else here",
            "Callers depend on this",
            "File: payment/PaymentMethod.java",
        ],
        size=15,
    )
    add_footer(slide)
    notes(
        slide,
        "PAYMENTMETHOD INTERFACE\n"
        "Objective: Read an interface as a SE contract, not a keyword demo.\n"
        "Why interfaces solve extensibility: new classes can satisfy the same promise.\n"
        "Analogy: A power socket standard — any compliant plug works; appliances vary.\n"
        "Live coding: Open PaymentMethod.java; walk Javadoc.\n"
        "Common mistakes: Putting System.out or card fields on the interface.\n"
        "Ask: Who benefits more — CardPayment author or Order.checkout author?\n"
        "Transition: Look at three implementations that honour the contract.",
    )


def slide_three_rails(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 3 · Implementations")
    content_title(slide, "Three workers. One job description.")

    rails = [
        ("CardPayment", "CARD", "Credit/debit simulation"),
        ("EftPayment", "EFT", "Bank transfer simulation"),
        ("WalletPayment", "WALLET", "In-app balance simulation"),
    ]
    for i, (cls, label, blurb) in enumerate(rails):
        left = Inches(0.7 + i * 4.1)
        add_round_rect(slide, left, Inches(2.0), Inches(3.9), Inches(4.3), WHITE)
        add_rect(slide, left, Inches(2.0), Inches(3.9), Inches(0.1), TEAL)
        add_textbox(
            slide,
            left + Inches(0.25),
            Inches(2.35),
            Inches(3.4),
            Inches(0.5),
            cls,
            size=16,
            bold=True,
            color=NAVY,
            font=FONT_TITLE,
        )
        add_round_rect(slide, left + Inches(0.25), Inches(3.1), Inches(2.2), Inches(0.55), TEAL)
        add_textbox(
            slide,
            left + Inches(0.25),
            Inches(3.18),
            Inches(2.2),
            Inches(0.4),
            label,
            size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.25),
            Inches(4.0),
            Inches(3.4),
            Inches(1.8),
            f"{blurb}\n\nimplements PaymentMethod\nprocessPayment → true\ngetPaymentType → \"{label}\"",
            size=14,
            color=SLATE,
        )
    add_footer(slide)
    notes(
        slide,
        "THREE RAILS\n"
        "Objective: See interchangeable implementations under one type.\n"
        "Live coding: Open CardPayment, EftPayment, WalletPayment side by side.\n"
        "Note: processPayment currently always returns true — simulation for Capstone 2.\n"
        "Ask: What would change inside Order if we only add WalletPayment? (Ideally: nothing.)\n"
        "Transition: Call sites must use the interface type — programming to abstractions / DIP.",
    )


def slide_program_to_abstractions(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 3 · DIP")
    content_title(
        slide,
        "Program to abstractions",
        "Dependency Inversion: high-level checkout depends on PaymentMethod, not CardPayment.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(5.9), Inches(4.3), CODE_BG)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Fragile",
        size=14,
        bold=True,
        color=AMBER,
        font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.7), Inches(5.3), Inches(3.3), "")
    add_multiline(
        tf,
        [
            "void checkout(CardPayment card) {",
            "  card.processPayment(total);",
            "}",
            "",
            "// Locked to one concrete class",
        ],
        size=14,
    )

    add_round_rect(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(4.3), CODE_BG)
    add_textbox(
        slide,
        Inches(7.2),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Stable",
        size=14,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(7.2), Inches(2.7), Inches(5.3), Inches(3.3), "")
    add_multiline(
        tf2,
        [
            "void checkout(PaymentMethod m) {",
            "  // any rail works",
            "  ...",
            "}",
            "",
            "// Depends on the contract",
        ],
        size=14,
    )
    add_footer(slide)
    notes(
        slide,
        "PROGRAM TO ABSTRACTIONS / DIP\n"
        "Objective: State DIP in Capstone 2 language without heavy theory.\n"
        "Why: Inverting dependencies protects stable business flow from volatile rail details.\n"
        "Analogy: Hire 'a driver' not 'specifically Thabo in a Toyota'.\n"
        "Live coding: Show Order.checkout(PaymentMethod) signature and PayNestApplication "
        "PaymentMethod paymentMethod = new CardPayment();\n"
        "Common mistakes: Parameter type CardPayment everywhere 'because the demo uses card'.\n"
        "Ask: Which version survives adding EftPayment without recompiling Order?\n"
        "Transition: A contract is useless unless polymorphism actually selects behaviour — Block 4.",
    )


# ---------------------------------------------------------------------------
# Block 4 — Polymorphism & SOLID
# ---------------------------------------------------------------------------


def slide_polymorphism(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 4 · Polymorphism")
    content_title(
        slide,
        "Same call site. Different behaviour.",
        "Polymorphism removes the if/switch from checkout.",
    )

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.5), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.1), Inches(2.3), Inches(11.2), Inches(3.9), "")
    add_multiline(
        tf,
        [
            "PaymentMethod paymentMethod = new CardPayment();",
            "order.checkout(paymentMethod);",
            "",
            "// Swap the rail — checkout code does not change:",
            "PaymentMethod paymentMethod = new EftPayment();",
            "order.checkout(paymentMethod);",
            "",
            "// Runtime picks CardPayment.processPayment or EftPayment.processPayment",
            "// That selection is dynamic dispatch.",
        ],
        size=16,
    )
    add_footer(slide)
    notes(
        slide,
        "POLYMORPHISM\n"
        "Objective: Explain polymorphism as one variable, many possible behaviours.\n"
        "Why: Eliminates type switches at the orchestration layer.\n"
        "Analogy: Pressing 'Play' on different media players — same button, different content.\n"
        "Live coding: In PayNestApplication, change CardPayment to EftPayment, rerun, compare console.\n"
        "Common mistakes: Still writing if (method.getPaymentType().equals(\"CARD\")) in checkout.\n"
        "Ask: Where did the else-if go?\n"
        "Transition: Name Open/Closed — we are open for new rails, closed for Order arithmetic edits.",
    )


def slide_dynamic_dispatch(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 4 · Dynamic dispatch")
    content_title(
        slide,
        "Dynamic dispatch (without the scary jargon)",
        "At runtime, Java asks the actual object: which processPayment should run?",
    )

    steps = [
        ("1", "Compile time", "Variable type is PaymentMethod — only interface methods allowed."),
        ("2", "Runtime object", "Heap holds a CardPayment or EftPayment instance."),
        ("3", "Dispatch", "method.processPayment(amount) jumps to that class's override."),
        ("4", "Result", "Checkout stays one path; behaviour varies by object."),
    ]
    for i, (n, h, b) in enumerate(steps):
        top = Inches(1.95 + i * 1.15)
        add_round_rect(slide, Inches(0.7), top, Inches(0.7), Inches(0.9), TEAL)
        add_textbox(
            slide,
            Inches(0.7),
            top + Inches(0.2),
            Inches(0.7),
            Inches(0.5),
            n,
            size=22,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font=FONT_TITLE,
        )
        add_textbox(
            slide,
            Inches(1.7),
            top + Inches(0.1),
            Inches(10.5),
            Inches(0.35),
            h,
            size=16,
            bold=True,
            color=NAVY,
            font=FONT_TITLE,
        )
        add_textbox(
            slide, Inches(1.7), top + Inches(0.45), Inches(10.5), Inches(0.4), b, size=15, color=SLATE
        )
    add_footer(slide)
    notes(
        slide,
        "DYNAMIC DISPATCH\n"
        "Objective: Demystify how polymorphism is executed.\n"
        "Keep JVM vtable talk light — second-year level.\n"
        "Ask: If the variable type is PaymentMethod, can checkout call a Card-only private method? (No.)\n"
        "Transition: This is what makes Open/Closed practical.",
    )


def slide_ocp(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 4 · Open/Closed")
    content_title(
        slide,
        "Open/Closed Principle emerges",
        "Open for extension. Closed for modification of Order's core arithmetic.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(5.9), Inches(4.3), TEAL)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.35),
        Inches(5.3),
        Inches(0.5),
        "Open for extension",
        size=20,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(3.2),
        Inches(5.3),
        Inches(2.5),
        "Add BuyNowPayLaterPayment\nas a new class that\nimplements PaymentMethod.\n\nWire it in the demo.",
        size=18,
        color=OFF_WHITE,
    )

    add_round_rect(slide, Inches(6.9), Inches(2.0), Inches(5.9), Inches(4.3), NAVY)
    add_textbox(
        slide,
        Inches(7.2),
        Inches(2.35),
        Inches(5.3),
        Inches(0.5),
        "Closed for modification",
        size=20,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(7.2),
        Inches(3.2),
        Inches(5.3),
        Inches(2.5),
        "Do not reopen\nOrder.calculateTotal()\nto teach a new rail.\n\nTotals stay Capstone 1 truth.",
        size=18,
        color=OFF_WHITE,
    )
    add_footer(slide)
    notes(
        slide,
        "OCP\n"
        "Objective: State OCP with Capstone 2 examples.\n"
        "Why: SOLID names patterns teams already earned through pain.\n"
        "Analogy: USB ports — new devices without redesigning the laptop motherboard.\n"
        "Distinction drill: What files change for BNPL? New class + wiring — not Order math.\n"
        "Ask: Does OCP mean we never edit any file? (No — we avoid editing the wrong files.)\n"
        "Transition: Who orchestrates talking to the rail — Order alone, or a collaborator?",
    )


def slide_exercise_files_change(prs):
    slide = exercise_slide(
        prs,
        "Distinction drill: BNPL file list",
        "With the interface design, list the files you expect to touch to add BuyNowPayLaterPayment.\n\n"
        "Which files should stay untouched?\n\n"
        "Compare to your Block 1 if/else list.",
        "5 min",
    )
    notes(
        slide,
        "EXERCISE — FILES CHANGE\n"
        "Good answer: new BuyNowPayLaterPayment.java; maybe PayNestApplication wiring; optional test.\n"
        "Untouched: Order.calculateTotal, OrderItem, Product, ideally Order.checkout body.\n"
        "Celebrate the shorter list vs Block 1.\n"
        "Transition: Block 5 — PaymentProcessor and checkout orchestration.",
    )


# ---------------------------------------------------------------------------
# Block 5 — Orchestration
# ---------------------------------------------------------------------------


def slide_processor(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 5 · PaymentProcessor")
    content_title(
        slide,
        "Delegation & separation of concerns",
        "PaymentProcessor talks to any PaymentMethod and owns the success/fail narrative.",
    )

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(12), Inches(4.5), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.25), Inches(11.5), Inches(4.0), "")
    add_multiline(
        tf,
        [
            "public void processPayment(PaymentMethod method, double amount) {",
            "  boolean success = method.processPayment(amount);",
            "  if (success) {",
            "    System.out.println(\"Payment successful via \" + method.getPaymentType());",
            "    System.out.println(\"Amount: R\" + String.format(\"%.0f\", amount));",
            "  } else {",
            "    System.out.println(\"Payment failed via \" + method.getPaymentType());",
            "  }",
            "}",
        ],
        size=15,
    )
    add_footer(slide)
    notes(
        slide,
        "PAYMENTPROCESSOR\n"
        "Objective: See SRP / SoC — processing messages ≠ owning line items.\n"
        "Why: Delegation lets Order request payment without owning rail messaging.\n"
        "Analogy: A restaurant runner takes the bill to the till; the chef does not swipe cards.\n"
        "Live coding: Open PaymentProcessor.java; trace method → processPayment → prints.\n"
        "Ask: Could Order print these lines itself? Yes — why don't we prefer that?\n"
        "Transition: Order.checkout wires total + processor + completion.",
    )


def slide_checkout(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 5 · Order.checkout")
    content_title(
        slide,
        "Checkout = total once, then delegate",
        "DRY: the amount charged comes from calculateTotal() — one source of truth.",
    )

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(7.5), Inches(4.5), CODE_BG)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.25), Inches(7.0), Inches(4.0), "")
    add_multiline(
        tf,
        [
            "public void checkout(PaymentMethod paymentMethod) {",
            "  double total = calculateTotal();",
            "  PaymentProcessor processor = new PaymentProcessor();",
            "  processor.processPayment(paymentMethod, total);",
            "  System.out.println(",
            "    \"Order completed successfully.\");",
            "}",
        ],
        size=14,
    )

    add_round_rect(slide, Inches(8.5), Inches(1.95), Inches(4.2), Inches(4.5), WHITE)
    add_textbox(
        slide,
        Inches(8.8),
        Inches(2.2),
        Inches(3.7),
        Inches(0.4),
        "Teaching note",
        size=14,
        bold=True,
        color=AMBER,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(8.8),
        Inches(2.8),
        Inches(3.7),
        Inches(3.2),
        "new PaymentProcessor() inside Order is a Capstone 2 trade-off.\n\n"
        "Later capstones push dependencies further outward (injection / routing).\n\n"
        "For now: clear demo + correct direction of abstraction.",
        size=14,
        color=SLATE,
    )
    add_footer(slide)
    notes(
        slide,
        "ORDER.CHECKOUT\n"
        "Objective: Trace the orchestration path end to end.\n"
        "Concepts: delegation, SoC, SRP, DRY, composition of payment.\n"
        "Live coding: Open Order.checkout; then PayNestApplication Capstone 2 block; run demo.\n"
        "Common mistakes: Recalculating total differently for payment; hard-coding amount 12400.\n"
        "Ask: If Mouse quantity changes, does checkout still charge the right amount? Why?\n"
        "Transition: Design without proof regresses — tests freeze the contract.",
    )


def slide_srp_soc(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 5 · SRP & SoC")
    content_title(slide, "Who owns what?")

    rows = [
        ("Order", "Lines, totals, summary, start checkout"),
        ("PaymentMethod", "Rail-specific charge + type label"),
        ("PaymentProcessor", "Invoke method + print success/fail story"),
        ("PayNestApplication", "Composition root / demo wiring"),
    ]
    add_round_rect(slide, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.55), TEAL)
    add_textbox(
        slide, Inches(1.0), Inches(2.0), Inches(3.5), Inches(0.35), "Type", size=13, bold=True, color=WHITE
    )
    add_textbox(
        slide,
        Inches(4.8),
        Inches(2.0),
        Inches(7.5),
        Inches(0.35),
        "Responsibility",
        size=13,
        bold=True,
        color=WHITE,
    )
    for i, (a, b) in enumerate(rows):
        top = Inches(2.55 + i * 0.95)
        bg = WHITE if i % 2 == 0 else ROW_ALT
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(0.85), bg)
        add_textbox(
            slide, Inches(1.0), top + Inches(0.2), Inches(3.5), Inches(0.45), a, size=16, bold=True, color=NAVY
        )
        add_textbox(slide, Inches(4.8), top + Inches(0.2), Inches(7.5), Inches(0.45), b, size=15, color=SLATE)
    add_footer(slide)
    notes(
        slide,
        "SRP / SoC TABLE\n"
        "Objective: Map classes to single clear jobs.\n"
        "Package boundaries: domain vs payment vs app.\n"
        "Transition: Business rules and tests.",
    )


# ---------------------------------------------------------------------------
# Block 6 — Rules & testing
# ---------------------------------------------------------------------------


def slide_business_rules(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 6 · Business rules")
    content_title(slide, "Rules the design must obey")

    rules = [
        ("Amount charged", "Equals calculateTotal() for that checkout — no silent discounts."),
        ("Stable labels", "CARD / EFT / WALLET — human-readable for logs and receipts."),
        ("One funnel", "Checkout accepts PaymentMethod — not a concrete rail type."),
        ("Extension test", "New rail = new class + wiring — not Order arithmetic edits."),
    ]
    for i, (h, b) in enumerate(rules):
        top = Inches(1.95 + i * 1.15)
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(1.0), WHITE)
        add_rect(slide, Inches(0.7), top, Inches(0.12), Inches(1.0), TEAL)
        add_textbox(
            slide,
            Inches(1.1),
            top + Inches(0.15),
            Inches(3.2),
            Inches(0.7),
            h,
            size=15,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            Inches(4.5),
            top + Inches(0.2),
            Inches(7.8),
            Inches(0.65),
            b,
            size=15,
            color=NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    add_footer(slide)
    notes(
        slide,
        "BUSINESS RULES\n"
        "Objective: Tie design choices to merchant-visible truth.\n"
        "Why: Architecture that violates business rules is still wrong.\n"
        "DRY: one total path shared by summary and checkout.\n"
        "Ask: Where would a silent discount bug hide?\n"
        "Transition: Tests prove rails and protect Capstone 1 totals.",
    )


def slide_testing(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 6 · Testing")
    content_title(
        slide,
        "Unit tests & regression",
        "Polymorphism is not proven until ≥2 rails appear in tests.",
    )

    add_round_rect(slide, Inches(0.7), Inches(1.95), Inches(5.9), Inches(4.4), WHITE)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Already green (C1)",
        size=16,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.8), Inches(5.3), Inches(3.2), "")
    add_bullets(
        tf,
        [
            "OrderTest totals & quantities",
            "Unmodifiable getItems",
            "OrderItem / OrderService tests",
            "Regression: must stay green after C2",
        ],
        size=15,
    )

    add_round_rect(slide, Inches(6.9), Inches(1.95), Inches(5.9), Inches(4.4), NAVY)
    add_textbox(
        slide,
        Inches(7.2),
        Inches(2.2),
        Inches(5.3),
        Inches(0.4),
        "Still required (C2)",
        size=16,
        bold=True,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(7.2), Inches(2.8), Inches(5.3), Inches(3.2), "")
    add_bullets(
        tf2,
        [
            "Checkout with CardPayment",
            "Checkout with a second rail",
            "Prove interface type is used",
            "mvn test passes on the branch",
        ],
        size=15,
        color=OFF_WHITE,
    )
    add_footer(slide)
    notes(
        slide,
        "TESTING\n"
        "Objective: Distinguish unit proof of polymorphism from C1 regression.\n"
        "Why: Without tests, a 'clean' interface can still be unused at call sites.\n"
        "Live coding: Open OrderTest; note absence of payment tests as a gap students fill.\n"
        "Common mistakes: Only testing CardPayment; testing getPaymentType in isolation only.\n"
        "Ask: Why is one rail in tests not enough for the rubric?\n"
        "Transition: Sketch a test exercise.",
    )


def slide_exercise_tests(prs):
    slide = exercise_slide(
        prs,
        "Exercise: sketch two-rail tests",
        "On paper, outline JUnit tests that:\n"
        "1) Build the laptop+mouse order (R12400)\n"
        "2) Checkout with CardPayment\n"
        "3) Checkout with WalletPayment (fresh order or reset)\n\n"
        "What would you assert? What would you deliberately NOT hard-code?",
        "6 min",
    )
    notes(
        slide,
        "EXERCISE — TESTS\n"
        "Guide: assert totals still 12400; assert getPaymentType labels; optionally capture stdout.\n"
        "Do not hard-code payment amount separate from calculateTotal().\n"
        "Transition: Markers score the same story — Block 7 rubric mapping.",
    )


# ---------------------------------------------------------------------------
# Block 7 — Rubric
# ---------------------------------------------------------------------------


def slide_rubric_map(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 7 · Rubric")
    content_title(slide, "How today's concepts become marks")

    rows = [
        ("Architecture & polymorphism", "25%", "Interface, DIP, OCP, no type switches"),
        ("Correctness & business rules", "30%", "Total → amount; labels; success path"),
        ("Testing & verification", "20%", "≥2 rails; C1 regression; mvn test"),
        ("Code quality & maintainability", "15%", "SRP, DRY, cohesion, packages"),
        ("Documentation & communication", "10%", "Why interfaces; new-rail answer"),
    ]
    add_round_rect(slide, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.5), TEAL)
    for label, left, w in [
        ("Category", Inches(0.9), Inches(4.5)),
        ("Wt", Inches(5.6), Inches(0.8)),
        ("Concepts from this lesson", Inches(6.6), Inches(5.7)),
    ]:
        add_textbox(slide, left, Inches(1.92), w, Inches(0.35), label, size=12, bold=True, color=WHITE)

    for i, (cat, wt, concepts) in enumerate(rows):
        top = Inches(2.45 + i * 0.8)
        bg = WHITE if i % 2 == 0 else ROW_ALT
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(0.72), bg)
        add_textbox(slide, Inches(0.9), top + Inches(0.15), Inches(4.5), Inches(0.45), cat, size=13, bold=True, color=NAVY)
        add_textbox(slide, Inches(5.6), top + Inches(0.15), Inches(0.8), Inches(0.45), wt, size=13, bold=True, color=TEAL)
        add_textbox(slide, Inches(6.6), top + Inches(0.15), Inches(5.7), Inches(0.45), concepts, size=13, color=SLATE)
    add_footer(slide)
    notes(
        slide,
        "RUBRIC MAP\n"
        "Objective: Students can explain how SE concepts map to Capstone 2 scoring.\n"
        "Walk each row; connect to a slide from earlier in the lesson.\n"
        "Ask: Which weight rewards 'it works in the demo' vs 'it can grow'?\n"
        "Transition: Distinction signal — timed BNPL.",
    )


def slide_distinction(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 7 · Distinction")
    content_title(
        slide,
        "Level 4 signal: timed BNPL extension",
        "New rail needs new class + wiring — not Order edits.",
    )

    add_round_rect(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.3), NAVY)
    add_textbox(
        slide,
        Inches(1.2),
        Inches(2.5),
        Inches(11),
        Inches(3.2),
        "Reviewers will ask:\n\n"
        "\"If we add BuyNowPayLaterPayment next sprint, what files change?\"\n\n"
        "Best answer: one new class + composition root / demo wiring.\n"
        "Weak answer: reopen Order and add another else-if.",
        size=20,
        color=OFF_WHITE,
    )
    add_footer(slide)
    notes(
        slide,
        "DISTINCTION\n"
        "Optional in-class sprint: 10 minutes to sketch BNPL class on paper.\n"
        "Transition: Closing takeaway.",
    )


def slide_concept_glossary(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Reference")
    content_title(slide, "Concept → Capstone 2 artefact")

    pairs = [
        ("Encapsulation / getters", "Order owns items; unmodifiable view"),
        ("Composition", "Order has items; checkout uses a PaymentMethod"),
        ("Interface / abstraction", "PaymentMethod"),
        ("Polymorphism / dispatch", "checkout(PaymentMethod) + rail overrides"),
        ("DIP", "Parameter type is the interface"),
        ("OCP", "New rail without editing totals"),
        ("SRP / SoC / delegation", "Order vs Processor vs Method"),
        ("DRY / business rules", "calculateTotal() feeds payment amount"),
        ("Unit + regression tests", "≥2 rails; C1 OrderTest stays green"),
    ]
    for i, (a, b) in enumerate(pairs):
        col, row = i % 2, i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.85 + row * 1.0)
        add_round_rect(slide, left, top, Inches(5.95), Inches(0.88), WHITE)
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.1),
            Inches(5.5),
            Inches(0.3),
            a,
            size=12,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
        )
        add_textbox(
            slide, left + Inches(0.2), top + Inches(0.42), Inches(5.5), Inches(0.35), b, size=13, color=NAVY
        )
    add_footer(slide)
    notes(slide, "GLOSSARY — leave up during Q&A or print as handout from speaker view.")


def slide_expected_output(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Block 5 · Demo")
    content_title(
        slide,
        "Expected console narrative",
        "Same Capstone 1 summary — then polymorphic payment confirmation.",
    )
    add_round_rect(slide, Inches(2.5), Inches(1.95), Inches(8.3), Inches(4.6), CODE_BG)
    _, tf = add_textbox(slide, Inches(2.9), Inches(2.25), Inches(7.6), Inches(4.1), "")
    add_multiline(
        tf,
        [
            "Order Summary",
            "Customer: John Smith",
            "Items:",
            "Laptop x1 - R12000",
            "Mouse x2 - R400",
            "Total: R12400",
            "",
            "Payment successful via CARD",
            "Amount: R12400",
            "Order completed successfully.",
        ],
        size=16,
    )
    add_footer(slide)
    notes(
        slide,
        "EXPECTED OUTPUT\n"
        "Live coding: run mvn exec:java; reconcile Mouse 200×2 and grand total by hand.\n"
        "Swap to EftPayment and show only the rail label changes.\n"
        "Ask: Which number must never silently diverge — summary total or payment amount?",
    )


def slide_common_mistakes(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Watch-outs")
    content_title(slide, "Common mistakes to catch early")

    mistakes = [
        ("Type switches return", "if (getPaymentType().equals(\"CARD\")) inside checkout"),
        ("Concrete parameters", "checkout(CardPayment) instead of PaymentMethod"),
        ("Hard-coded amount", "Charging 12400 instead of calculateTotal()"),
        ("Mutable guts", "Exposing the live ArrayList of items"),
        ("One-rail tests", "Only CardPayment — polymorphism never proven"),
        ("Docs skip the why", "No answer to 'what changes for a new rail?'"),
    ]
    for i, (h, b) in enumerate(mistakes):
        col, row = i % 2, i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.95 + row * 1.5)
        add_round_rect(slide, left, top, Inches(5.95), Inches(1.3), WHITE)
        add_rect(slide, left, top, Inches(0.1), Inches(1.3), RED_SOFT)
        add_textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.2),
            Inches(5.4),
            Inches(0.35),
            h,
            size=15,
            bold=True,
            color=NAVY,
            font=FONT_TITLE,
        )
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.65), Inches(5.4), Inches(0.45), b, size=14, color=SLATE
        )
    add_footer(slide)
    notes(
        slide,
        "COMMON MISTAKES\n"
        "Use mid-lesson or as a closing checklist before students start Capstone 2 alone.\n"
        "Ask students to mark which mistake they are most tempted by.",
    )


def slide_questions_bank(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Facilitation")
    content_title(slide, "Questions to keep asking")

    qs = [
        "What must NOT change when Ops adds a new payment rail?",
        "Is Order's job totals — or Visa vs EFT protocols?",
        "Who benefits more from PaymentMethod: rail authors or checkout?",
        "Where did the else-if go after polymorphism?",
        "What files change for BuyNowPayLaterPayment?",
        "Why is one rail in tests not enough for the rubric?",
    ]
    for i, q in enumerate(qs):
        top = Inches(1.9 + i * 0.75)
        add_round_rect(slide, Inches(0.7), top, Inches(0.55), Inches(0.55), TEAL)
        add_textbox(
            slide,
            Inches(0.7),
            top + Inches(0.08),
            Inches(0.55),
            Inches(0.4),
            str(i + 1),
            size=16,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font=FONT_TITLE,
        )
        add_textbox(slide, Inches(1.5), top + Inches(0.1), Inches(11), Inches(0.45), q, size=16, color=NAVY)
    add_footer(slide)
    notes(
        slide,
        "QUESTIONS BANK\n"
        "Park this slide for cold-call moments; answers appear in earlier speaker notes.",
    )


def slide_closing(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.8),
        Inches(11.5),
        Inches(1.4),
        "Interfaces exist because\nsoftware changes.",
        size=36,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1.5),
        Inches(3.6),
        Inches(10.3),
        Inches(0.8),
        "Capstone 2 is the design you earned by feeling the if/else pain.",
        size=18,
        color=TEAL_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_round_rect(slide, Inches(2.8), Inches(4.7), Inches(7.7), Inches(1.4), TEAL)
    add_textbox(
        slide,
        Inches(3.0),
        Inches(4.95),
        Inches(7.3),
        Inches(1.0),
        "Homework brief:\ndocs/assessments/capstone-02-payment-methods.md",
        size=16,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    notes(
        slide,
        "CLOSING\n"
        "Point to assessment brief and optional DDD live session for deeper language.\n"
        "Remind: mvn test; demo with polymorphic checkout; short design note for docs marks.\n"
        "Thank students.",
    )


def slide_file_open_order(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Lecturer prep")
    content_title(slide, "Suggested file open order", "Walk-only — narrate existing solution.")

    files = [
        "1. PayNestApplication.java — C1 then C2 blocks",
        "2. Order.java — totals, encapsulation, checkout",
        "3. (Scratch) anti-pattern if/else — do not commit",
        "4. PaymentMethod.java — contract reveal",
        "5. CardPayment / EftPayment / WalletPayment",
        "6. PaymentProcessor.java",
        "7. OrderTest.java — regression + C2 gap",
        "8. Capstone 2 brief + rubric CSV",
    ]
    for i, line in enumerate(files):
        top = Inches(1.9 + i * 0.58)
        add_textbox(slide, Inches(0.9), top, Inches(11.5), Inches(0.5), line, size=16, color=NAVY)
    add_footer(slide)
    notes(slide, "PREP SLIDE — for lecturer; skip or show briefly at start.")


def build():
    prs = new_presentation()

    # Opening
    slide_title(prs)
    slide_agenda(prs)
    slide_not_syntax(prs)
    section_divider(
        prs,
        "0",
        "Capstone 1 recap",
        "~15 minutes",
        "Reconnect objects, encapsulation, and trustworthy totals before payment pressure arrives.",
    )
    slide_you_already_know(prs)
    slide_c1_walk(prs)
    slide_hook(prs)

    # Block 1
    section_divider(
        prs,
        "1",
        "Why software changes",
        "~25 minutes",
        "Explain change pressure and why tightly coupled checkout is hard to maintain.",
    )
    slide_why_software_changes(prs)
    slide_bad_checkout(prs)
    slide_exercise_bnpl_pain(prs)
    slide_tight_coupling(prs)

    # Block 2
    section_divider(
        prs,
        "2",
        "Responsibility & structure",
        "~20 minutes",
        "Name responsibility, cohesion, coupling, and composition using Capstone 1 code.",
    )
    slide_responsibility(prs)
    slide_cohesion_coupling(prs)
    slide_composition(prs)
    slide_exercise_vocab(prs)

    # Block 3
    section_divider(
        prs,
        "3",
        "Abstractions & interfaces",
        "~25 minutes",
        "Derive a shared payment contract and program to that abstraction.",
    )
    slide_need_abstraction(prs)
    slide_exercise_signatures(prs)
    slide_payment_method(prs)
    slide_three_rails(prs)
    slide_program_to_abstractions(prs)

    # Block 4
    section_divider(
        prs,
        "4",
        "Polymorphism & SOLID",
        "~30 minutes",
        "Use polymorphism and dynamic dispatch to remove switches; recognise OCP and DIP.",
    )
    slide_polymorphism(prs)
    slide_dynamic_dispatch(prs)
    slide_ocp(prs)
    slide_exercise_files_change(prs)

    # Block 5
    section_divider(
        prs,
        "5",
        "Checkout orchestration",
        "~20 minutes",
        "Trace PaymentProcessor and Order.checkout as separation of concerns and delegation.",
    )
    slide_processor(prs)
    slide_checkout(prs)
    slide_expected_output(prs)
    slide_srp_soc(prs)

    # Block 6
    section_divider(
        prs,
        "6",
        "Rules, DRY & testing",
        "~20 minutes",
        "Lock business rules with unit tests across ≥2 rails and Capstone 1 regression.",
    )
    slide_business_rules(prs)
    slide_testing(prs)
    slide_exercise_tests(prs)

    # Block 7
    section_divider(
        prs,
        "7",
        "Rubric mapping",
        "~15 minutes",
        "Map every lesson concept onto Capstone 2 assessment weights.",
    )
    slide_rubric_map(prs)
    slide_distinction(prs)
    slide_concept_glossary(prs)
    slide_common_mistakes(prs)
    slide_questions_bank(prs)
    slide_file_open_order(prs)
    slide_closing(prs)

    DOCS_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    prs.save(DOCS_OUT)
    print(f"Saved {OUT_FILE} ({len(prs.slides)} slides)")
    print(f"Saved {DOCS_OUT}")


if __name__ == "__main__":
    build()
