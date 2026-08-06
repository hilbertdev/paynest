#!/usr/bin/env python3
"""Build student-facing OOP Foundations PPTX (no speaker notes)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
DOCS_OUT = (
    Path(__file__).resolve().parents[2]
    / "docs"
    / "live-sessions"
    / "oop-foundations-capstone-02"
    / "oop-foundations-capstone-02.pptx"
)
OUT_FILE = OUT_DIR / "oop-foundations-capstone-02.pptx"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED_SOFT = RGBColor(0xB9, 0x1C, 0x1C)
GREEN = RGBColor(0x05, 0x96, 0x69)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
ROW_ALT = RGBColor(0xF1, 0xF5, 0xF9)

FONT_TITLE = "Poppins"
FONT_BODY = "Lato"
FONT_CODE = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER = "PayNest · OOP Foundations · Capstone 2 Prep"


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box, tf


def add_bullets(tf, items, size=17, color=SLATE, spacing=8):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)


def add_code_lines(tf, lines, size=13, color=OFF_WHITE, spacing=2):
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_CODE
        p.space_after = Pt(spacing)


def add_header_bar(slide, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL)
    if kicker:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(0.28),
            Inches(12),
            Inches(0.32),
            kicker.upper(),
            size=11,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
        )


def add_footer(slide, text=FOOTER):
    add_textbox(
        slide, Inches(0.7), Inches(7.05), Inches(11.5), Inches(0.3), text, size=9, color=SLATE
    )


def content_title(slide, title, top=0.65):
    add_textbox(
        slide,
        Inches(0.7),
        Inches(top),
        Inches(12),
        Inches(0.65),
        title,
        size=28,
        bold=True,
        color=NAVY,
        font=FONT_TITLE,
    )


def discussion(slide, text, top=6.35):
    add_round_rect(slide, Inches(0.7), Inches(top), Inches(11.9), Inches(0.55), WHITE)
    add_rect(slide, Inches(0.7), Inches(top), Inches(0.1), Inches(0.55), AMBER)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(top + 0.08),
        Inches(11.3),
        Inches(0.4),
        f"Discussion: {text}",
        size=13,
        color=SLATE,
    )


def standard_slide(prs, kicker, title):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, kicker)
    content_title(slide, title)
    add_footer(slide)
    return slide


def add_table(slide, left, top, width, height, headers, rows):
    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height)
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.name = FONT_TITLE
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = FONT_BODY
                p.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else ROW_ALT
    return table


def code_panel(slide, left, top, width, height, lines, size=13):
    add_round_rect(slide, left, top, width, height, CODE_BG)
    _, tf = add_textbox(
        slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.35), height - Inches(0.25), ""
    )
    add_code_lines(tf, lines, size=size)


def pill(slide, left, top, width, height, text, fill=TEAL, color=WHITE):
    add_round_rect(slide, left, top, width, height, fill)
    add_textbox(
        slide,
        left,
        top + Inches(0.12),
        width,
        height - Inches(0.1),
        text,
        size=13,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        font=FONT_TITLE,
    )


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def slide_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(11.5),
        Inches(1.2),
        "Object-Oriented Programming Foundations",
        size=36,
        bold=True,
        color=WHITE,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(3.3),
        Inches(11),
        Inches(0.5),
        "Preparing for Capstone 2",
        size=22,
        color=TEAL_LIGHT,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(4.2),
        Inches(11),
        Inches(1.0),
        "PayNest checkout · interfaces · polymorphism\n"
        "Extend payment rails without rewriting order logic",
        size=16,
        color=OFF_WHITE,
    )
    discussion(slide, "What broke last time when one change forced edits in three files?", top=5.8)


def slide_outcomes(prs):
    slide = standard_slide(prs, "Part 1 · Introduction", "Learning Outcomes")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(4.2), "")
    add_bullets(
        tf,
        [
            "Name the four pillars of OOP and why they exist",
            "Explain encapsulation using PayNest’s Order / Product",
            "Explain polymorphism using PaymentMethod",
            "See how polymorphism leads into SOLID (OCP, DIP)",
            "Sketch the Capstone 2 architecture before you write it",
        ],
        size=20,
        spacing=14,
    )
    discussion(slide, "Which outcome feels most unfamiliar right now?")


def slide_why_oop(prs):
    slide = standard_slide(prs, "Part 1 · Introduction", "Why Object-Oriented Programming?")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.45), Inches(6.2), Inches(3.8), "")
    add_bullets(
        tf,
        [
            "Real-world modelling — orders, customers, payments → objects",
            "Reusable code — calculateTotal() written once",
            "Easier maintenance — fix line rules in OrderItem",
            "Extensibility — add Wallet later without rewriting arithmetic",
        ],
        size=16,
        spacing=12,
    )
    add_round_rect(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(3.6), WHITE)
    add_rect(slide, Inches(7.2), Inches(1.5), Inches(0.12), Inches(3.6), TEAL)
    add_textbox(
        slide,
        Inches(7.55),
        Inches(1.7),
        Inches(4.7),
        Inches(0.4),
        "Analogy",
        size=14,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    add_textbox(
        slide,
        Inches(7.55),
        Inches(2.2),
        Inches(4.7),
        Inches(2.6),
        "A Nest thermostat: you press “heat.” "
        "You do not rewire the furnace each time.\n\n"
        "Stable contract · replaceable implementation.",
        size=15,
        color=SLATE,
    )
    discussion(slide, "In Capstone 1, what stayed stable when Mouse quantity became 2?")


def slide_pillars(prs):
    slide = standard_slide(prs, "Part 2 · The Four Pillars", "The Four Pillars (Preview)")
    add_table(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(11.9),
        Inches(3.6),
        ["Pillar", "One-line meaning", "PayNest hint"],
        [
            ["Encapsulation", "Hide internals; expose a safe API", "Order private items"],
            ["Abstraction", "Show the essential contract", "PaymentMethod"],
            ["Inheritance", "Share / specialise behaviour", "Later: BasePaymentProvider"],
            ["Polymorphism", "One type, many behaviours", "checkout(PaymentMethod)"],
        ],
    )
    discussion(slide, "Which pillar sounds like “a menu every rail must support”?")


def slide_encap_what(prs):
    slide = standard_slide(prs, "Part 3 · Encapsulation", "Encapsulation — What Is It?")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.4), Inches(5.8), Inches(3.2), "")
    add_bullets(
        tf,
        [
            "Bundle data + behaviour, and control access",
            "Fields hold state (price, items, quantity)",
            "Methods define allowed operations",
            "Callers use the public API, not the raw guts",
        ],
        size=16,
        spacing=10,
    )
    code_panel(
        slide,
        Inches(6.8),
        Inches(1.4),
        Inches(5.8),
        Inches(3.5),
        [
            "public class Product {",
            "  private final int id;",
            "  private final String name;",
            "  private final double price;",
            "",
            "  public double getPrice() {",
            "    return price;",
            "  }",
            "}",
        ],
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.1),
        Inches(11.5),
        Inches(0.4),
        "Why should we care? So nobody sets price = -999 mid-checkout.",
        size=15,
        bold=True,
        color=NAVY,
    )
    discussion(slide, "Should a receipt printer change Product.price?")


def slide_private(prs):
    slide = standard_slide(prs, "Part 3 · Encapsulation", "Why private Fields Exist")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(12),
        Inches(2.2),
        [
            "// Order.java",
            "private final int id;",
            "private final Customer customer;",
            "private final List<OrderItem> items;  // owned here",
        ],
        size=15,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(3.9), Inches(11.5), Inches(2.0), "")
    add_bullets(
        tf,
        [
            "Callers cannot order.items.clear()",
            "They must use addItem(product, quantity)",
            "Totals stay trustworthy — one mutation path",
            "Why care? Broken totals = angry merchants and failed demos",
        ],
        size=16,
        spacing=8,
    )
    discussion(slide, "What goes wrong if items is public?")


def slide_bad_encap(prs):
    slide = standard_slide(prs, "Part 3 · Encapsulation", "BAD Encapsulation")
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(11),
        Inches(0.35),
        "Fake anti-pattern (not in PayNest — on purpose)",
        size=14,
        bold=True,
        color=RED_SOFT,
    )
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.75),
        Inches(7.2),
        Inches(3.3),
        [
            "public class BadOrder {",
            "  public List<OrderItem> items = new ArrayList<>();",
            "  public double total;  // anyone can set this!",
            "",
            "  public void printSummary() {",
            "    System.out.println(\"Total: R\" + total);",
            "  }",
            "}",
        ],
    )
    _, tf = add_textbox(slide, Inches(8.2), Inches(1.75), Inches(4.4), Inches(3.3), "")
    add_bullets(
        tf,
        [
            "Caller can clear() items after totals",
            "Caller can set total = 0",
            "No single source of truth",
            "Bugs hide in silent mutation",
        ],
        size=15,
        spacing=10,
    )
    discussion(slide, "Worse: a wrong total, or a total that looks right but was hand-edited?")


def slide_good_encap(prs):
    slide = standard_slide(prs, "Part 3 · Encapsulation", "GOOD Encapsulation (PayNest Order)")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(3.5),
        [
            "public void addItem(Product product, int quantity) {",
            "  OrderItem orderItem = new OrderItem(product, quantity);",
            "  items.add(orderItem);",
            "}",
            "",
            "public List<OrderItem> getItems() {",
            "  return Collections.unmodifiableList(items);",
            "}",
            "",
            "public double calculateTotal() { /* sum line totals */ }",
        ],
        size=14,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.1),
        Inches(11.5),
        Inches(0.5),
        "Mutation only via addItem · read-only view outward · total is computed — never a free-floating field",
        size=14,
        color=GREEN,
    )
    discussion(slide, "Why unmodifiableList instead of the raw ArrayList?")


def slide_encap_quiz(prs):
    slide = standard_slide(prs, "Part 3 · Encapsulation", "Quiz — Encapsulation")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(4.0), "")
    add_bullets(
        tf,
        [
            "1. Why are Product fields private final?",
            "2. What happens if a caller add()s to order.getItems()?",
            "3. Where is quantity validated — Order or OrderItem?",
            "4. True/False: Getters always break encapsulation.",
        ],
        size=20,
        spacing=16,
    )
    discussion(slide, "Encapsulation = “never share data,” or “share it safely”?")


def slide_transition(prs):
    slide = standard_slide(prs, "Part 4 · Transition", "From Protection to Behaviour")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(1.5), "")
    add_bullets(
        tf,
        [
            "Encapsulation protects an object’s state",
            "Capstone 2 asks something new…",
        ],
        size=18,
        spacing=10,
    )
    add_round_rect(slide, Inches(0.7), Inches(3.2), Inches(11.9), Inches(2.2), WHITE)
    add_rect(slide, Inches(0.7), Inches(3.2), Inches(0.12), Inches(2.2), TEAL)
    add_textbox(
        slide,
        Inches(1.1),
        Inches(3.5),
        Inches(11),
        Inches(1.6),
        "How do Card, EFT, and Wallet behave differently\n"
        "while checkout stays one reusable path?\n\n"
        "We need a contract every rail obeys — then swap the rail.",
        size=18,
        color=NAVY,
    )
    discussion(slide, "If you wrote if (type.equals(\"CARD\"))…, what happens when BNPL arrives?")


def slide_poly_def(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Polymorphism — Definition")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.4), Inches(6.0), Inches(3.5), "")
    add_bullets(
        tf,
        [
            "Polymorphism = many forms",
            "One interface (shared contract)",
            "Many implementations (different behaviour)",
            "Callers program against the contract, not the concrete class",
        ],
        size=17,
        spacing=12,
    )
    code_panel(
        slide,
        Inches(6.9),
        Inches(1.4),
        Inches(5.7),
        Inches(3.5),
        [
            "PaymentMethod paymentMethod",
            "    = new CardPayment();",
            "order.checkout(paymentMethod);",
            "",
            "// Same checkout works for",
            "// EftPayment or WalletPayment",
        ],
        size=14,
    )
    discussion(slide, "What type is the variable? What type is the object?")


def slide_compile_runtime(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Compile-Time vs Runtime")
    add_table(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(11.9),
        Inches(2.4),
        ["Kind", "When decided", "Capstone 2?"],
        [
            ["Compile-time", "Before run (e.g. overloading)", "Rarely the focus"],
            ["Runtime", "During run (interface / override)", "This is Capstone 2"],
        ],
    )
    code_panel(
        slide,
        Inches(0.7),
        Inches(4.1),
        Inches(11.9),
        Inches(1.7),
        [
            "method.processPayment(amount);",
            "// Compiler: method is a PaymentMethod",
            "// JVM: picks Card / EFT / Wallet override",
        ],
        size=14,
    )
    discussion(slide, "If the variable is PaymentMethod, how does Java print \"CARD\"?")


def slide_interfaces(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Interfaces — The Contract")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(2.8),
        [
            "// com.paynestsystem.payment.PaymentMethod",
            "public interface PaymentMethod {",
            "    boolean processPayment(double amount);",
            "    String getPaymentType();",
            "}",
        ],
        size=16,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(4.5), Inches(11.5), Inches(1.5), "")
    add_bullets(
        tf,
        [
            "Abstraction: process this Rand amount; name your rail",
            "Small on purpose — easy to implement",
            "Checkout depends on this, not Card/EFT/Wallet",
        ],
        size=16,
        spacing=8,
    )
    discussion(slide, "Why two methods instead of one mega-method?")


def slide_implementations(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Implementations — Many Forms")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(7.0),
        Inches(3.6),
        [
            "public class CardPayment",
            "        implements PaymentMethod {",
            "  @Override",
            "  public boolean processPayment(double amount) {",
            "    return true; // simulated",
            "  }",
            "  @Override",
            "  public String getPaymentType() {",
            "    return \"CARD\";",
            "  }",
            "}",
        ],
        size=13,
    )
    # hierarchy cards
    pill(slide, Inches(8.2), Inches(1.5), Inches(4.3), Inches(0.7), "«interface» PaymentMethod", NAVY)
    pill(slide, Inches(8.2), Inches(2.6), Inches(4.3), Inches(0.55), "CardPayment", TEAL)
    pill(slide, Inches(8.2), Inches(3.35), Inches(4.3), Inches(0.55), "EftPayment", TEAL)
    pill(slide, Inches(8.2), Inches(4.1), Inches(4.3), Inches(0.55), "WalletPayment", TEAL)
    add_textbox(
        slide,
        Inches(8.2),
        Inches(4.85),
        Inches(4.3),
        Inches(0.7),
        "Same shape · different labels\nNo extends between rails",
        size=13,
        color=SLATE,
    )
    discussion(slide, "Related by extends? Why or why not?")


def slide_dispatch(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Dynamic Dispatch")
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(11),
        Inches(0.4),
        "JVM chooses the override at runtime.",
        size=16,
        color=SLATE,
    )
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.85),
        Inches(12),
        Inches(3.2),
        [
            "// PaymentProcessor.java",
            "public void processPayment(PaymentMethod method, double amount) {",
            "    boolean success = method.processPayment(amount);",
            "    if (success) {",
            "        System.out.println(\"Payment successful via \"",
            "            + method.getPaymentType());",
            "    }",
            "}",
        ],
        size=14,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.25),
        Inches(11.5),
        Inches(0.4),
        "Never mentions CardPayment · swap the object → swap the behaviour",
        size=15,
        bold=True,
        color=NAVY,
    )
    discussion(slide, "Where would an if/else anti-pattern have lived instead?")


def slide_checkout_needs(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Why Checkout Only Needs PaymentMethod")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(7.2),
        Inches(3.4),
        [
            "public void checkout(PaymentMethod paymentMethod) {",
            "  double total = calculateTotal();",
            "  PaymentProcessor processor =",
            "      new PaymentProcessor();",
            "  processor.processPayment(paymentMethod, total);",
            "}",
        ],
        size=13,
    )
    # flow boxes
    y = 1.5
    for i, (label, fill) in enumerate(
        [
            ("Order", NAVY),
            ("PaymentProcessor", TEAL),
            ("PaymentMethod", CODE_BG),
            ("Card / EFT / Wallet", TEAL_LIGHT),
        ]
    ):
        pill(slide, Inches(8.3), Inches(y + i * 0.85), Inches(4.2), Inches(0.6), label, fill)
    discussion(slide, "What changes when we add BuyNowPayLaterPayment?")


def slide_poly_action(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Polymorphism in Action")
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(11.9),
        Inches(1.8),
        [
            "// PayNestApplication — Capstone 2 block",
            "PaymentMethod paymentMethod = new CardPayment();",
            "order.checkout(paymentMethod);",
        ],
        size=15,
    )
    add_table(
        slide,
        Inches(0.7),
        Inches(3.4),
        Inches(11.9),
        Inches(2.2),
        ["Construction", "Type label", "Checkout edits?"],
        [
            ["new CardPayment()", "CARD", "No"],
            ["new EftPayment()", "EFT", "No"],
            ["new WalletPayment()", "WALLET", "No"],
        ],
    )
    discussion(slide, "Is polymorphism in the construction line or the checkout line?", top=6.4)


def slide_poly_quiz(prs):
    slide = standard_slide(prs, "Part 5 · Polymorphism", "Quick Check — Polymorphism")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(4.0), "")
    add_bullets(
        tf,
        [
            "1. What keyword links CardPayment to PaymentMethod?",
            "2. Can PaymentProcessor compile if CryptoPayment does not exist yet?",
            "3. True/False: Polymorphism requires extends.",
            "4. Name the runtime mechanism that picks the override.",
        ],
        size=19,
        spacing=14,
    )
    discussion(slide, "Why is “False” on Q3 important for Capstone 2?")


def slide_mechanism(prs):
    slide = standard_slide(prs, "Part 6 · Why This Matters", "Mechanism vs Design Principle")
    add_table(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(11.9),
        Inches(2.8),
        ["", "Mechanism", "Design principle"],
        [
            ["What", "Language feature", "Guidance for structure"],
            ["Example", "Polymorphism, interfaces", "SOLID (OCP, DIP, …)"],
            ["Role", "How behaviour varies", "How to organise modules"],
        ],
    )
    add_round_rect(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(1.2), WHITE)
    add_rect(slide, Inches(0.7), Inches(4.6), Inches(0.12), Inches(1.2), AMBER)
    add_textbox(
        slide,
        Inches(1.1),
        Inches(4.85),
        Inches(11),
        Inches(0.8),
        "Polymorphism is NOT a SOLID principle.\n"
        "It is the mechanism that enables several SOLID principles.",
        size=16,
        bold=True,
        color=NAVY,
    )
    discussion(slide, "Can you follow OCP without polymorphism? How painful?")


def slide_ocp(prs):
    slide = standard_slide(prs, "Part 6 · Why This Matters", "Open/Closed Principle (OCP)")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.4), Inches(6.2), Inches(3.5), "")
    add_bullets(
        tf,
        [
            "Open for extension · Closed for modification",
            "Add BuyNowPayLaterPayment → new class",
            "Do not rewrite calculateTotal() or processor logic",
            "Polymorphism is how OCP becomes real in Java",
        ],
        size=17,
        spacing=12,
    )
    add_round_rect(slide, Inches(7.2), Inches(1.45), Inches(5.3), Inches(3.6), WHITE)
    add_textbox(
        slide,
        Inches(7.5),
        Inches(1.65),
        Inches(4.8),
        Inches(0.4),
        "Closed",
        size=13,
        bold=True,
        color=SLATE,
        font=FONT_TITLE,
    )
    pill(slide, Inches(7.5), Inches(2.15), Inches(4.7), Inches(0.5), "Order", NAVY)
    pill(slide, Inches(7.5), Inches(2.8), Inches(4.7), Inches(0.5), "PaymentProcessor", NAVY)
    add_textbox(
        slide,
        Inches(7.5),
        Inches(3.5),
        Inches(4.8),
        Inches(0.35),
        "Open",
        size=13,
        bold=True,
        color=TEAL,
        font=FONT_TITLE,
    )
    pill(slide, Inches(7.5), Inches(3.95), Inches(4.7), Inches(0.5), "+ BuyNowPayLaterPayment", TEAL)
    discussion(slide, "Which files must change for BNPL? Which must not?")


def slide_dip(prs):
    slide = standard_slide(prs, "Part 6 · Why This Matters", "Dependency Inversion (DIP Preview)")
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(11),
        Inches(0.4),
        "Depend on abstractions, not concretions.",
        size=17,
        bold=True,
        color=NAVY,
    )
    code_panel(
        slide,
        Inches(0.7),
        Inches(1.9),
        Inches(12),
        Inches(2.8),
        [
            "// Good — Capstone 2",
            "public void checkout(PaymentMethod paymentMethod) { ... }",
            "",
            "// Weaker — honest look at PayNest today",
            "PaymentProcessor processor = new PaymentProcessor();",
        ],
        size=14,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.0),
        Inches(11.5),
        Inches(0.5),
        "Celebrate the PaymentMethod dependency; later capstones push more behind interfaces.",
        size=14,
        color=SLATE,
    )
    discussion(slide, "Why depend on PaymentMethod instead of CardPayment?")


def slide_poly_solid(prs):
    slide = standard_slide(prs, "Part 6 · Why This Matters", "Polymorphism Enables SOLID")
    cards = [
        ("OCP", "Open/Closed\nnew rail = new class"),
        ("DIP", "Dependency Inversion\ndepend on PaymentMethod"),
        ("LSP", "Liskov Substitution\nany rail works in checkout"),
        ("SRP", "Single Responsibility\nOrder totals vs rail details"),
    ]
    for i, (title, body) in enumerate(cards):
        left = Inches(0.7 + i * 3.1)
        add_round_rect(slide, left, Inches(1.55), Inches(2.9), Inches(3.2), WHITE)
        add_rect(slide, left, Inches(1.55), Inches(2.9), Inches(0.12), TEAL)
        add_textbox(
            slide,
            left + Inches(0.15),
            Inches(1.9),
            Inches(2.6),
            Inches(0.5),
            title,
            size=22,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.15),
            Inches(2.6),
            Inches(2.6),
            Inches(1.8),
            body,
            size=14,
            color=SLATE,
            align=PP_ALIGN.CENTER,
        )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.0),
        Inches(11.5),
        Inches(0.5),
        "Polymorphism ≠ SOLID · Capstone 2 grades architecture & polymorphism heavily (25%)",
        size=14,
        color=NAVY,
    )
    discussion(slide, "Which SOLID letter matches “add a rail without rewriting order logic”?")


def slide_architecture(prs):
    slide = standard_slide(prs, "Part 7 · Capstone Preview", "Capstone 2 Architecture")
    labels = [
        ("PayNestApplication", NAVY, 0.7),
        ("Order", TEAL, 3.5),
        ("PaymentProcessor", TEAL, 6.3),
        ("PaymentMethod", CODE_BG, 9.1),
    ]
    for text, fill, left in labels:
        pill(slide, Inches(left), Inches(1.6), Inches(2.5), Inches(0.7), text, fill)
    for left in (3.2, 6.0, 8.8):
        add_textbox(
            slide,
            Inches(left),
            Inches(1.7),
            Inches(0.4),
            Inches(0.5),
            "→",
            size=22,
            bold=True,
            color=SLATE,
            align=PP_ALIGN.CENTER,
        )
    pill(slide, Inches(2.5), Inches(3.0), Inches(2.5), Inches(0.55), "CardPayment", TEAL_LIGHT)
    pill(slide, Inches(5.4), Inches(3.0), Inches(2.5), Inches(0.55), "EftPayment", TEAL_LIGHT)
    pill(slide, Inches(8.3), Inches(3.0), Inches(2.5), Inches(0.55), "WalletPayment", TEAL_LIGHT)
    _, tf = add_textbox(slide, Inches(0.7), Inches(4.0), Inches(11.5), Inches(1.8), "")
    add_bullets(
        tf,
        [
            "Order — lines, totals, starts checkout",
            "PaymentProcessor — talks to any rail",
            "PaymentMethod — the contract",
            "Card / EFT / Wallet — concrete rails",
        ],
        size=16,
        spacing=8,
    )
    discussion(slide, "Where does the charged amount come from — rail or order?")


def slide_extending(prs):
    slide = standard_slide(prs, "Part 7 · Capstone Preview", "Extending Without Rewriting")
    add_round_rect(slide, Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.1), WHITE)
    add_rect(slide, Inches(0.7), Inches(1.4), Inches(0.12), Inches(1.1), TEAL)
    add_textbox(
        slide,
        Inches(1.1),
        Inches(1.65),
        Inches(11),
        Inches(0.7),
        "New rail = new class + wiring,\nnot editing Order’s core arithmetic.",
        size=18,
        bold=True,
        color=NAVY,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(2.8), Inches(11.5), Inches(2.8), "")
    add_bullets(
        tf,
        [
            "1. BuyNowPayLaterPayment implements PaymentMethod",
            "2. Wire in PayNestApplication (or a test)",
            "3. Leave calculateTotal() alone",
            "4. Leave PaymentProcessor alone (if it only uses the interface)",
            "That is OCP powered by polymorphism.",
        ],
        size=17,
        spacing=10,
    )
    discussion(slide, "Would a giant switch (type) inside Order pass distinction?")


def slide_remember(prs):
    slide = standard_slide(prs, "Part 7 · Capstone Preview", "What You Should Remember")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.45), Inches(11.5), Inches(3.5), "")
    add_bullets(
        tf,
        [
            "Encapsulation protects totals (private, addItem, unmodifiable views)",
            "Interfaces define payment contracts without payment details",
            "Polymorphism keeps checkout stable while rails vary",
            "OCP / DIP = design goals; polymorphism = Java mechanism",
            "Success test: new rail, minimal edits",
        ],
        size=17,
        spacing=12,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.2),
        Inches(11.5),
        Inches(0.5),
        "Same order total everywhere; only the payment rail changes.",
        size=16,
        bold=True,
        color=TEAL,
    )
    discussion(slide, "Explain Capstone 2 in one sentence to a classmate who missed today.")


def slide_live_coding(prs):
    slide = standard_slide(prs, "Appendix A", "Live Coding Plan (20 min)")
    add_table(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(11.9),
        Inches(4.6),
        ["Min", "Activity", "File(s)"],
        [
            ["0–3", "Run C1 summary + C2 checkout", "PayNestApplication"],
            ["3–7", "Break encapsulation (getItems().add)", "Order"],
            ["7–11", "if/else anti-pattern (scratch only)", "(do not commit)"],
            ["11–15", "PaymentMethod + three rails", "payment/*"],
            ["15–18", "Swap Card → EFT", "PayNestApplication"],
            ["18–20", "Sketch BNPL; list files touched", "whiteboard"],
        ],
    )


def slide_exercises(prs):
    slide = standard_slide(prs, "Appendix B", "Student Exercises")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(4.8), "")
    add_bullets(
        tf,
        [
            "Trace the call stack from main to CardPayment.processPayment",
            "Predict the exception from order.getItems().clear()",
            "List three if/else payment designs and why each fails OCP",
            "Pair: add a console line in each rail naming rail + amount",
            "Test sketch: JUnit outline for two rails’ getPaymentType()",
        ],
        size=18,
        spacing=14,
    )


def slide_misconceptions(prs):
    slide = standard_slide(prs, "Appendix C", "Common Misconceptions")
    add_table(
        slide,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(5.0),
        ["Misconception", "Reality"],
        [
            ["Interfaces are “enterprise only”", "Capstone 2 is 1 interface + 3 small classes"],
            ["Polymorphism requires extends", "implements is enough — preferred here"],
            ["Getters break encapsulation", "Uncontrolled mutation breaks it"],
            ["Polymorphism is a SOLID principle", "It is a mechanism; SOLID are principles"],
            ["Order should know card fees", "Rails own rail details; Order owns totals"],
            ["BNPL means editing PaymentProcessor", "Only if you hard-coded card logic"],
        ],
    )


def slide_quiz(prs):
    slide = standard_slide(prs, "Appendix D", "Quiz Questions")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.45), Inches(11.5), Inches(5.0), "")
    add_bullets(
        tf,
        [
            "1. Quote the two methods on PaymentMethod.",
            "2. Why does Order.getItems() use unmodifiableList?",
            "3. In PaymentMethod m = new WalletPayment(); — variable type? object type?",
            "4. Dynamic dispatch in one sentence (PaymentProcessor).",
            "5. One OCP-friendly change and one OCP-breaking change for a new rail.",
            "6. Is new PaymentProcessor() inside Order ideal DIP? Why/why not?",
            "7. Capstone rubric question about BNPL — what do reviewers ask?",
        ],
        size=16,
        spacing=10,
    )


def slide_homework(prs):
    slide = standard_slide(prs, "Appendix E", "Homework (before Capstone 2)")
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.45), Inches(11.5), Inches(4.5), "")
    add_bullets(
        tf,
        [
            "Read docs/assessments/capstone-02-payment-methods.md end-to-end",
            "Draw: PaymentMethod ← Card / EFT / Wallet",
            "One paragraph: why interfaces beat a mega-method",
            "Run mvn test and mvn exec:java; note the payment lines",
            "Stretch: draft BuyNowPayLaterPayment locally",
            "Stretch: sketch a test that checkouts with two different rails",
            "Bring: diagram + “files changed for BNPL” list",
        ],
        size=16,
        spacing=10,
    )


def slide_refs(prs):
    slide = standard_slide(prs, "References", "PayNest Source Paths")
    add_table(
        slide,
        Inches(0.5),
        Inches(1.45),
        Inches(12.3),
        Inches(4.8),
        ["Topic", "Path"],
        [
            ["Demo wiring", "app/PayNestApplication.java"],
            ["Encapsulation", "domain/Order.java, OrderItem, Product"],
            ["Contract", "payment/PaymentMethod.java"],
            ["Rails", "CardPayment, EftPayment, WalletPayment"],
            ["Orchestration", "payment/PaymentProcessor.java"],
            ["Brief", "docs/assessments/capstone-02-payment-methods.md"],
        ],
    )


def build():
    prs = new_presentation()
    builders = [
        slide_title,
        slide_outcomes,
        slide_why_oop,
        slide_pillars,
        slide_encap_what,
        slide_private,
        slide_bad_encap,
        slide_good_encap,
        slide_encap_quiz,
        slide_transition,
        slide_poly_def,
        slide_compile_runtime,
        slide_interfaces,
        slide_implementations,
        slide_dispatch,
        slide_checkout_needs,
        slide_poly_action,
        slide_poly_quiz,
        slide_mechanism,
        slide_ocp,
        slide_dip,
        slide_poly_solid,
        slide_architecture,
        slide_extending,
        slide_remember,
        slide_live_coding,
        slide_exercises,
        slide_misconceptions,
        slide_quiz,
        slide_homework,
        slide_refs,
    ]
    for fn in builders:
        fn(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    DOCS_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    prs.save(str(DOCS_OUT))
    print(f"Wrote {OUT_FILE}")
    print(f"Wrote {DOCS_OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
